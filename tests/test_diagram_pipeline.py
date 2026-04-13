"""Tests for the diagram native rendering pipeline integration."""

import tempfile
from pathlib import Path

from pptx import Presentation

from edupptx.diagram_native import SlotPosition, draw_diagram_on_slide
from edupptx.layout_resolver import resolve_layout, _compute_material_slot
from edupptx.models import (
    ContentMaterial, PresentationPlan, SlideCard, SlideContent, ResolvedSlide,
)
from edupptx.pipeline_v2 import render_with_schema
from edupptx.style_resolver import resolve_style
from edupptx.style_schema import load_style


STYLES_DIR = Path(__file__).parent.parent / "styles"


def _emerald_style():
    return resolve_style(load_style(STYLES_DIR / "emerald.json"))


def _flowchart_data():
    return {
        "nodes": [
            {"id": "1", "label": "Start"},
            {"id": "2", "label": "Process"},
            {"id": "3", "label": "End"},
        ],
        "edges": [
            {"from": "1", "to": "2"},
            {"from": "2", "to": "3"},
        ],
        "direction": "TB",
    }


def _plan_with_diagram(diagram_type: str, diagram_data: dict) -> PresentationPlan:
    return PresentationPlan(
        topic="Test",
        palette="emerald",
        slides=[
            SlideContent(
                type="content",
                title="Slide with diagram",
                cards=[SlideCard(icon="star", title="Point", body="Body")],
                notes="Notes",
                content_materials=[
                    ContentMaterial(
                        action="generate_diagram",
                        position="center",
                        diagram_type=diagram_type,
                        diagram_data=diagram_data,
                    )
                ],
            ),
        ],
    )


# ── resolve_layout generates diagram_info ──


def test_resolve_layout_with_diagram_specs():
    """resolve_layout should populate diagram_info when diagram_specs provided."""
    plan = _plan_with_diagram("flowchart", _flowchart_data())
    style = _emerald_style()
    diagram_specs = {0: ("flowchart", _flowchart_data())}

    slides = resolve_layout(plan, style, diagram_specs=diagram_specs)

    assert len(slides) == 1
    info = slides[0].diagram_info
    assert info is not None
    d_type, d_data, slot = info
    assert d_type == "flowchart"
    assert d_data == _flowchart_data()
    assert len(slot) == 4
    assert all(isinstance(v, int) for v in slot)


def test_resolve_layout_without_diagram_no_info():
    """resolve_layout without diagram_specs should leave diagram_info as None."""
    plan = _plan_with_diagram("flowchart", _flowchart_data())
    style = _emerald_style()

    slides = resolve_layout(plan, style)

    assert len(slides) == 1
    assert slides[0].diagram_info is None


def test_resolve_layout_diagram_no_image_shape():
    """Diagram slides should NOT have an image shape (diagrams are native vectors)."""
    plan = _plan_with_diagram("flowchart", _flowchart_data())
    style = _emerald_style()
    diagram_specs = {0: ("flowchart", _flowchart_data())}

    slides = resolve_layout(plan, style, diagram_specs=diagram_specs)

    image_shapes = [s for s in slides[0].shapes if s.shape_type == "image"]
    assert len(image_shapes) == 0


# ── draw_diagram_on_slide for all 5 types ──


def _make_blank_slide():
    """Create a blank pptx slide for testing."""
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _test_slot():
    return SlotPosition(x=500000, y=500000, width=8000000, height=4000000)


def test_draw_flowchart():
    slide = _make_blank_slide()
    style = _emerald_style()
    draw_diagram_on_slide(slide, "flowchart", _flowchart_data(), _test_slot(), style)
    assert len(slide.shapes) > 0


def test_draw_timeline():
    slide = _make_blank_slide()
    style = _emerald_style()
    data = {
        "events": [
            {"year": "2020", "label": "Event A", "description": "Desc A"},
            {"year": "2021", "label": "Event B", "description": "Desc B"},
        ]
    }
    draw_diagram_on_slide(slide, "timeline", data, _test_slot(), style)
    assert len(slide.shapes) > 0


def test_draw_comparison():
    slide = _make_blank_slide()
    style = _emerald_style()
    data = {
        "columns": [
            {"header": "Pros", "items": ["Fast", "Simple"]},
            {"header": "Cons", "items": ["Limited"]},
        ]
    }
    draw_diagram_on_slide(slide, "comparison", data, _test_slot(), style)
    assert len(slide.shapes) > 0


def test_draw_hierarchy():
    slide = _make_blank_slide()
    style = _emerald_style()
    data = {
        "root": {
            "label": "Root",
            "children": [
                {"label": "Child A", "children": []},
                {"label": "Child B", "children": [
                    {"label": "Grandchild", "children": []}
                ]},
            ],
        }
    }
    draw_diagram_on_slide(slide, "hierarchy", data, _test_slot(), style)
    assert len(slide.shapes) > 0


def test_draw_cycle():
    slide = _make_blank_slide()
    style = _emerald_style()
    data = {
        "steps": [
            {"label": "Step 1", "description": "Desc 1"},
            {"label": "Step 2", "description": "Desc 2"},
            {"label": "Step 3", "description": "Desc 3"},
        ]
    }
    draw_diagram_on_slide(slide, "cycle", data, _test_slot(), style)
    assert len(slide.shapes) > 0


def test_draw_unknown_type_noop():
    """Unknown diagram type should not crash, just no shapes added."""
    slide = _make_blank_slide()
    style = _emerald_style()
    initial_count = len(slide.shapes)
    draw_diagram_on_slide(slide, "nonexistent", {}, _test_slot(), style)
    assert len(slide.shapes) == initial_count


# ── End-to-end: render_with_schema with diagram_specs ──


def test_render_with_diagram_flowchart():
    """Full pipeline: plan with flowchart diagram renders to valid PPTX."""
    plan = _plan_with_diagram("flowchart", _flowchart_data())
    diagram_specs = {0: ("flowchart", _flowchart_data())}

    with tempfile.TemporaryDirectory() as tmpdir:
        out = Path(tmpdir) / "test_diagram.pptx"
        render_with_schema(
            plan, STYLES_DIR / "emerald.json",
            diagram_specs=diagram_specs,
            output_path=out,
        )
        assert out.exists()
        prs = Presentation(str(out))
        assert len(prs.slides) == 1
        # Diagram should have added native shapes
        assert len(prs.slides[0].shapes) > 0


def test_render_with_all_diagram_types():
    """All 5 diagram types render without error in the full pipeline."""
    diagrams = {
        "flowchart": _flowchart_data(),
        "timeline": {"events": [{"year": "2020", "label": "E", "description": "D"}]},
        "comparison": {"columns": [{"header": "A", "items": ["x"]}]},
        "hierarchy": {"root": {"label": "R", "children": [{"label": "C", "children": []}]}},
        "cycle": {"steps": [{"label": "S1", "description": "D1"}, {"label": "S2", "description": "D2"}]},
    }

    for dtype, ddata in diagrams.items():
        plan = _plan_with_diagram(dtype, ddata)
        specs = {0: (dtype, ddata)}

        with tempfile.TemporaryDirectory() as tmpdir:
            out = Path(tmpdir) / f"test_{dtype}.pptx"
            render_with_schema(
                plan, STYLES_DIR / "emerald.json",
                diagram_specs=specs,
                output_path=out,
            )
            assert out.exists(), f"Failed for diagram type: {dtype}"


# ── Regression: material slot respects slide type layout ──


def test_material_slot_image_left_uses_left_half():
    """image_left slide: material slot must be in left half regardless of position."""
    style = _emerald_style()
    slot = _compute_material_slot("image_left", style, "center")
    assert slot is not None
    x, y, w, h = slot
    # Material should occupy left 50% of content area
    assert x == style.margin_left
    assert w == int(style.content_w * 0.50)
    # Should NOT be full width
    assert w < style.content_w


def test_material_slot_image_right_uses_right_half():
    """image_right slide: material slot must be in right half."""
    style = _emerald_style()
    slot = _compute_material_slot("image_right", style, "center")
    assert slot is not None
    x, y, w, h = slot
    # Material should be in right half
    mat_w = int(style.content_w * 0.50)
    assert w == mat_w
    assert x > style.margin_left  # offset to the right


def test_hierarchy_fits_within_slot():
    """Hierarchy diagram shapes must not exceed slot bottom boundary."""
    from pptx.util import Emu
    slide = _make_blank_slide()
    style = _emerald_style()
    # Small slot: 200pt tall (3-level tree needs ~150pt with default sizes)
    small_slot = SlotPosition(x=500000, y=500000, width=8000000, height=2540000)
    data = {
        "root": {
            "label": "Root",
            "children": [
                {"label": "A", "children": [
                    {"label": "A1", "children": []},
                    {"label": "A2", "children": []},
                ]},
                {"label": "B", "children": [
                    {"label": "B1", "children": []},
                ]},
            ],
        }
    }
    draw_diagram_on_slide(slide, "hierarchy", data, small_slot, style)
    # All shapes must be within slot bounds
    slot_bottom = small_slot.y + small_slot.height
    for shape in slide.shapes:
        shape_bottom = shape.top + shape.height
        assert shape_bottom <= slot_bottom + 12700, (  # 1pt tolerance
            f"Shape at y={shape.top} h={shape.height} exceeds slot bottom {slot_bottom}"
        )


def test_flowchart_tb_fits_within_slot():
    """TB flowchart with many nodes must scale to fit slot height."""
    slide = _make_blank_slide()
    style = _emerald_style()
    # Tiny slot: only 100pt tall, 6 nodes
    tiny_slot = SlotPosition(x=500000, y=500000, width=4000000, height=1270000)
    data = {
        "nodes": [{"id": str(i), "label": f"Step {i}"} for i in range(6)],
        "edges": [{"from": str(i), "to": str(i+1)} for i in range(5)],
        "direction": "TB",
    }
    draw_diagram_on_slide(slide, "flowchart", data, tiny_slot, style)
    slot_bottom = tiny_slot.y + tiny_slot.height
    for shape in slide.shapes:
        shape_bottom = shape.top + shape.height
        assert shape_bottom <= slot_bottom + 12700, (
            f"Shape at y={shape.top} h={shape.height} exceeds slot bottom {slot_bottom}"
        )
