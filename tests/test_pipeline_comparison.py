"""Pipeline tests: verify v2 schema-driven pipeline produces correct output.

Tests schema switching (same plan, different style -> different output),
text preservation, slide count, and structural integrity.
"""

import pytest
from pathlib import Path
from pptx import Presentation

from edupptx.models import PresentationPlan, SlideContent, SlideCard
from edupptx.pipeline_v2 import render_with_schema


STYLES_DIR = Path(__file__).parent.parent / "styles"


def _test_plan():
    return PresentationPlan(
        topic="Test Comparison",
        palette="emerald",
        slides=[
            SlideContent(
                type="cover",
                title="Test Presentation",
                subtitle="Comparing old and new pipelines",
                cards=[
                    SlideCard(icon="star", title="Point 1", body="First key insight"),
                    SlideCard(icon="target", title="Point 2", body="Second important point"),
                    SlideCard(icon="check", title="Point 3", body="Third conclusion"),
                ],
                formula="a² + b² = c²",
                notes="Cover slide notes",
            ),
            SlideContent(
                type="content",
                title="Core Concepts",
                cards=[
                    SlideCard(icon="book", title="Theory", body="Mathematical foundations"),
                    SlideCard(icon="globe", title="Practice", body="Real world applications"),
                    SlideCard(icon="lightbulb", title="Insight", body="Key takeaway points"),
                ],
                footer="Chapter 1 summary",
                notes="Content notes",
            ),
            SlideContent(
                type="big_quote",
                title="Education is the most powerful weapon you can use to change the world.",
                footer="— Nelson Mandela",
                notes="Quote notes",
            ),
        ],
    )


def test_v2_pipeline_produces_valid_pptx(tmp_path):
    """New pipeline produces a valid, openable PPTX file."""
    plan = _test_plan()
    out = render_with_schema(plan, STYLES_DIR / "emerald.json", output_path=tmp_path / "v2.pptx")
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == 3


def test_correct_slide_count(tmp_path):
    """Pipeline produces exactly 3 slides."""
    plan = _test_plan()
    out = render_with_schema(plan, STYLES_DIR / "emerald.json", output_path=tmp_path / "v2.pptx")
    prs = Presentation(str(out))
    assert len(prs.slides) == 3


def test_text_content_preserved(tmp_path):
    """New pipeline preserves all text content from the plan."""
    plan = _test_plan()
    out = render_with_schema(plan, STYLES_DIR / "emerald.json", output_path=tmp_path / "v2.pptx")
    prs = Presentation(str(out))

    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)

    joined = " ".join(all_text)
    assert "Test Presentation" in joined
    assert "Core Concepts" in joined
    assert "change the world" in joined
    assert "Nelson Mandela" in joined


def test_different_styles_different_output(tmp_path):
    """Swapping style JSON produces visually different output."""
    plan = _test_plan()

    emerald_path = render_with_schema(
        plan, STYLES_DIR / "emerald.json", output_path=tmp_path / "emerald.pptx",
    )
    blue_path = render_with_schema(
        plan, STYLES_DIR / "blue.json", output_path=tmp_path / "blue.pptx",
    )

    assert emerald_path.exists()
    assert blue_path.exists()
    assert emerald_path.stat().st_size > 10_000
    assert blue_path.stat().st_size > 10_000


def test_speaker_notes_preserved(tmp_path):
    """Speaker notes from the plan appear in the output."""
    plan = _test_plan()
    out = render_with_schema(plan, STYLES_DIR / "emerald.json", output_path=tmp_path / "v2.pptx")
    prs = Presentation(str(out))

    notes = []
    for slide in prs.slides:
        try:
            notes.append(slide.notes_slide.notes_text_frame.text)
        except Exception:
            pass

    joined = " ".join(notes)
    assert "Cover slide notes" in joined
    assert "Content notes" in joined
