"""Tests for the PPTX writer — validates output file structure."""

import pytest
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

from edupptx.models import (
    ResolvedFont, ResolvedShape, ResolvedShadow, ResolvedSlide,
)
from edupptx.pptx_writer import PptxWriter
from edupptx.style_schema import SLIDE_W, SLIDE_H


def _make_simple_slides():
    return [
        ResolvedSlide(
            shapes=[
                ResolvedShape(
                    shape_type="textbox",
                    left=100000, top=100000,
                    width=5000000, height=500000,
                    text="Hello World",
                    font=ResolvedFont("Noto Sans SC", "Arial", 24, True, "#1F2937"),
                    z_order=1,
                ),
                ResolvedShape(
                    shape_type="rounded_rect",
                    left=100000, top=700000,
                    width=3000000, height=2000000,
                    fill_color="#FFFFFF",
                    corner_radius=8000,
                    shadow=ResolvedShadow(381000, 101600, "#93C5FD", 14),
                    z_order=0,
                ),
            ],
            notes="Test notes",
        ),
    ]


def test_writer_creates_valid_pptx(tmp_path):
    writer = PptxWriter()
    writer.write_slides(_make_simple_slides())
    out = writer.save(tmp_path / "test.pptx")
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_writer_correct_dimensions(tmp_path):
    writer = PptxWriter()
    writer.write_slides(_make_simple_slides())
    out = writer.save(tmp_path / "test.pptx")
    prs = Presentation(str(out))
    assert prs.slide_width == Emu(SLIDE_W)
    assert prs.slide_height == Emu(SLIDE_H)


def test_writer_preserves_notes(tmp_path):
    writer = PptxWriter()
    writer.write_slides(_make_simple_slides())
    out = writer.save(tmp_path / "test.pptx")
    prs = Presentation(str(out))
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert "Test notes" in notes


def test_writer_multiple_slides(tmp_path):
    slides = [
        ResolvedSlide(shapes=[
            ResolvedShape(shape_type="textbox", left=0, top=0, width=1000000, height=500000,
                          text=f"Slide {i}", font=ResolvedFont("Arial", "Arial", 20), z_order=0),
        ], notes=f"Notes {i}")
        for i in range(5)
    ]
    writer = PptxWriter()
    writer.write_slides(slides)
    out = writer.save(tmp_path / "multi.pptx")
    prs = Presentation(str(out))
    assert len(prs.slides) == 5
