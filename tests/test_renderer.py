"""Tests for the v2 rendering pipeline (replaces old renderer tests)."""

import tempfile
from pathlib import Path

from pptx import Presentation

from edupptx.models import PresentationPlan, SlideCard, SlideContent, ContentMaterial
from edupptx.pipeline_v2 import render_with_schema
from edupptx.backgrounds import generate_background
from edupptx.style_resolver import resolve_style
from edupptx.style_schema import load_style


STYLES_DIR = Path(__file__).parent.parent / "styles"


def _make_simple_plan() -> PresentationPlan:
    return PresentationPlan(
        topic="Test",
        palette="emerald",
        slides=[
            SlideContent(
                type="cover",
                title="Test Presentation",
                subtitle="A test subtitle",
                cards=[
                    SlideCard(icon="star", title="Point 1", body="Description 1"),
                    SlideCard(icon="target", title="Point 2", body="Description 2"),
                    SlideCard(icon="check", title="Point 3", body="Description 3"),
                ],
                formula="E = mc²",
                notes="Test speaker notes",
            ),
            SlideContent(
                type="content",
                title="Content Slide",
                cards=[
                    SlideCard(icon="book", title="Topic A", body="Details about A"),
                    SlideCard(icon="lightbulb", title="Topic B", body="Details about B"),
                ],
                notes="More notes",
            ),
            SlideContent(
                type="closing",
                title="Thank You",
                subtitle="End of presentation",
                notes="Closing remarks",
            ),
        ],
    )


def test_renderer_creates_valid_pptx():
    """Render a simple plan and verify the output is a valid PPTX."""
    plan = _make_simple_plan()
    resolved = resolve_style(load_style(STYLES_DIR / "emerald.json"))

    with tempfile.TemporaryDirectory() as tmpdir:
        cache_dir = Path(tmpdir) / "cache"
        styles = ["diagonal_gradient", "radial_gradient", "geometric_circles"]
        backgrounds = [
            generate_background(resolved, styles[i % len(styles)], cache_dir)
            for i in range(len(plan.slides))
        ]

        out_path = Path(tmpdir) / "test.pptx"
        render_with_schema(plan, STYLES_DIR / "emerald.json",
                           bg_paths=backgrounds, output_path=out_path)

        assert out_path.exists()
        assert out_path.stat().st_size > 10000

        prs = Presentation(str(out_path))
        assert len(prs.slides) == 3


def test_renderer_slide_dimensions():
    """Verify slide dimensions match 16:9 standard."""
    plan = _make_simple_plan()
    with tempfile.TemporaryDirectory() as tmpdir:
        out_path = Path(tmpdir) / "test.pptx"
        render_with_schema(plan, STYLES_DIR / "emerald.json", output_path=out_path)
        prs = Presentation(str(out_path))
        assert prs.slide_width == 12192000
        assert prs.slide_height == 6858000


def test_renderer_speaker_notes():
    """Verify speaker notes are embedded."""
    plan = _make_simple_plan()
    resolved = resolve_style(load_style(STYLES_DIR / "emerald.json"))

    with tempfile.TemporaryDirectory() as tmpdir:
        cache_dir = Path(tmpdir) / "cache"
        styles = ["diagonal_gradient", "radial_gradient", "geometric_circles"]
        backgrounds = [
            generate_background(resolved, styles[i % len(styles)], cache_dir)
            for i in range(len(plan.slides))
        ]

        out_path = Path(tmpdir) / "test.pptx"
        render_with_schema(plan, STYLES_DIR / "emerald.json",
                           bg_paths=backgrounds, output_path=out_path)

        prs = Presentation(str(out_path))
        slide = prs.slides[0]
        assert slide.has_notes_slide
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Test speaker notes" in notes_text


def test_renderer_with_illustration():
    """Test rendering a slide with an illustration material."""
    from PIL import Image
    img = Image.new("RGB", (1024, 768), (200, 220, 200))
    img_path = Path(tempfile.mktemp(suffix=".png"))
    img.save(img_path, "PNG")

    plan = PresentationPlan(
        topic="Test",
        palette="emerald",
        slides=[
            SlideContent(
                type="content",
                title="Test with illustration",
                cards=[SlideCard(icon="star", title="Point 1", body="Body 1")],
                notes="Notes",
                content_materials=[
                    ContentMaterial(
                        action="generate_illustration",
                        position="center",
                        illustration_description="test illustration",
                    )
                ],
            ),
        ],
    )

    with tempfile.TemporaryDirectory() as tmpdir:
        out = Path(tmpdir) / "test.pptx"
        render_with_schema(
            plan, STYLES_DIR / "emerald.json",
            material_paths={0: img_path},
            output_path=out,
        )
        assert out.exists()
        prs = Presentation(str(out))
        assert len(prs.slides) == 1

    img_path.unlink(missing_ok=True)


def test_renderer_long_card_body():
    """Verify that long card body text doesn't crash the renderer."""
    plan = PresentationPlan(
        topic="Test",
        palette="emerald",
        slides=[
            SlideContent(
                type="content",
                title="Test",
                cards=[SlideCard(icon="star", title="Title",
                                 body="这是一段超过五十个字的卡片正文内容，用来测试当文本过长时渲染器是否能自动缩小字号而不报错。这段文字故意写得很长很长。")],
                notes="Notes",
            ),
        ],
    )
    with tempfile.TemporaryDirectory() as tmpdir:
        out = Path(tmpdir) / "test.pptx"
        render_with_schema(plan, STYLES_DIR / "emerald.json", output_path=out)
        assert out.exists()
