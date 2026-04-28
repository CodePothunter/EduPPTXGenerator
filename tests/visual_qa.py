"""Visual QA pipeline for generated PPTX files.

Inspects shape geometry to detect:
  - non-background overlaps > threshold
  - text frames likely to overflow their containers
  - shapes extending beyond slide bounds
  - slides that are mostly empty

Usage:
    python -m tests.visual_qa <pptx_path> [--out report.json]
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import asdict, dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# Decorative atmospheric shapes (large opacity ≤ 0.25 circles or rects) are
# intentionally allowed to bleed off-canvas and overlap content, so the QA
# checker excludes them from the bounds and overlap signals.
_DECORATIVE_ALPHA_THRESHOLD = 0.25
_ALPHA_RE = re.compile(r'<a:alpha\s+val="(\d+)"')


@dataclass
class Issue:
    type: str
    severity: str  # critical | high | medium | low
    description: str
    shapes_involved: list[str] = field(default_factory=list)


@dataclass
class SlideReport:
    slide_number: int
    issues: list[Issue] = field(default_factory=list)


def _shape_box(sh) -> tuple[int, int, int, int] | None:
    try:
        if sh.left is None or sh.top is None or sh.width is None or sh.height is None:
            return None
        return int(sh.left), int(sh.top), int(sh.width), int(sh.height)
    except Exception:
        return None


def _overlap_ratio(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> float:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0, min(ay + ah, by + bh) - max(ay, by))
    inter = ix * iy
    smaller = min(aw * ah, bw * bh)
    return inter / smaller if smaller > 0 else 0.0


def _is_contained(inner: tuple[int, int, int, int], outer: tuple[int, int, int, int],
                   slack: int = 50_000) -> bool:
    """True if `inner` sits (mostly) inside `outer` with EMU slack."""
    ix, iy, iw, ih = inner
    ox, oy, ow, oh = outer
    return (
        ix >= ox - slack
        and iy >= oy - slack
        and ix + iw <= ox + ow + slack
        and iy + ih <= oy + oh + slack
    )


def _shape_kind(sh) -> str:
    """Classify shape: 'text' if it has a text frame with content, else 'visual'."""
    if getattr(sh, "has_text_frame", False):
        try:
            if sh.text_frame.text.strip():
                return "text"
        except Exception:
            pass
    return "visual"


def _is_background(sh) -> bool:
    """Treat the slide-sized picture (the bg image) as background."""
    box = _shape_box(sh)
    if box is None:
        return False
    _, _, w, h = box
    # >85% of slide size is background-ish
    return w > Emu(9_000_000) and h > Emu(5_000_000)


def _shape_min_alpha(sh) -> float:
    """Return the smallest fill/stroke alpha (0..1). 1.0 if no alpha set."""
    try:
        xml = sh._element.xml
    except Exception:
        return 1.0
    matches = _ALPHA_RE.findall(xml)
    if not matches:
        return 1.0
    try:
        return min(int(m) / 100_000 for m in matches)
    except ValueError:
        return 1.0


def _is_decorative(sh) -> bool:
    """Low-opacity atmospheric shapes (decorative blobs) are intentionally
    allowed to bleed off-canvas — exclude them from QA signals."""
    return _shape_min_alpha(sh) <= _DECORATIVE_ALPHA_THRESHOLD


def _estimate_text_overflow(sh) -> bool:
    """Crude estimate: total chars * avg char height vs shape height."""
    if not sh.has_text_frame or sh.height in (None, 0):
        return False
    tf = sh.text_frame
    text = tf.text or ""
    if not text:
        return False
    line_count = max(1, text.count("\n") + 1)
    # Default font size estimate (pts → EMUs roughly)
    # 18pt font = ~228600 EMU line height
    line_h_emu = 228_600
    needed = line_count * line_h_emu
    return needed > int(sh.height) * 1.4  # 40% slack


def analyze_pptx(pptx_path: Path, *, overlap_threshold: float = 0.10) -> dict:
    pres = Presentation(str(pptx_path))
    slide_w = int(pres.slide_width)
    slide_h = int(pres.slide_height)

    reports: list[SlideReport] = []
    for idx, slide in enumerate(pres.slides, 1):
        rep = SlideReport(slide_number=idx)
        boxes: list[tuple[str, tuple[int, int, int, int], str]] = []
        non_bg_area = 0

        for sh in slide.shapes:
            box = _shape_box(sh)
            if box is None:
                continue
            name = sh.name or type(sh).__name__
            x, y, w, h = box
            decorative = _is_decorative(sh)

            if (x + w > slide_w or y + h > slide_h or x < 0 or y < 0) and not decorative:
                rep.issues.append(Issue(
                    "out_of_bounds", "high",
                    f"Shape '{name}' extends outside slide bounds",
                    [name],
                ))

            if _is_background(sh) or decorative:
                continue

            non_bg_area += w * h
            boxes.append((name, box, _shape_kind(sh)))

            if _estimate_text_overflow(sh):
                rep.issues.append(Issue(
                    "text_overflow", "medium",
                    f"Text in '{name}' likely overflows container",
                    [name],
                ))

        # pairwise overlap — only flag same-kind collisions (text↔text, visual↔visual)
        # and skip pairs where one is fully contained in the other (intentional layout).
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a_name, a_box, a_kind = boxes[i]
                b_name, b_box, b_kind = boxes[j]
                if a_kind != b_kind:
                    continue
                if _is_contained(a_box, b_box) or _is_contained(b_box, a_box):
                    continue
                ratio = _overlap_ratio(a_box, b_box)
                if ratio > overlap_threshold:
                    rep.issues.append(Issue(
                        "overlap", "medium" if ratio < 0.4 else "high",
                        f"{a_kind.title()} shapes overlap by {ratio:.0%}",
                        [a_name, b_name],
                    ))

        # emptiness
        slide_area = slide_w * slide_h
        if slide_area > 0:
            density = non_bg_area / slide_area
            if density < 0.30:
                rep.issues.append(Issue(
                    "sparse_slide", "low",
                    f"Slide content covers only {density:.0%} of area",
                ))

        reports.append(rep)

    # summary
    severity_counts: dict[str, int] = {}
    for rep in reports:
        for iss in rep.issues:
            severity_counts[iss.severity] = severity_counts.get(iss.severity, 0) + 1

    return {
        "pptx": str(pptx_path),
        "slide_count": len(reports),
        "severity_counts": severity_counts,
        "critical_issues": severity_counts.get("critical", 0),
        "high_issues": severity_counts.get("high", 0),
        "slides": [
            {"slide_number": r.slide_number, "issues": [asdict(i) for i in r.issues]}
            for r in reports
        ],
    }


def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser(description="Visual QA for generated PPTX")
    p.add_argument("pptx", type=Path)
    p.add_argument("--out", type=Path, default=None, help="Write JSON report")
    p.add_argument("--quiet", action="store_true", help="Only print summary line")
    args = p.parse_args(argv)

    if not args.pptx.exists():
        print(f"PPTX not found: {args.pptx}", file=sys.stderr)
        return 2

    report = analyze_pptx(args.pptx)
    if args.out:
        args.out.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")

    summary = (
        f"slides={report['slide_count']} "
        f"critical={report['critical_issues']} "
        f"high={report['high_issues']} "
        f"counts={report['severity_counts']}"
    )
    if args.quiet:
        print(summary)
    else:
        print(json.dumps(report, ensure_ascii=False, indent=2))
        print(summary, file=sys.stderr)

    return 1 if report["critical_issues"] > 0 else 0


if __name__ == "__main__":
    sys.exit(main())
