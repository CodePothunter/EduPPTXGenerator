"""Assemble SVG files into a PPTX presentation.

Three modes:
- Default: SVG→native DrawingML shapes (directly editable, no "Convert to Shape")
- Embed: SVG+PNG dual embedding (ZIP post-processing, image-only, --embed flag)
- Legacy native: svg2pptx library (--legacy-native flag)

The default mode parses each SVG element and converts it to a native PowerPoint
shape (rect, text box, path, etc.), producing fully editable output.
"""

from __future__ import annotations

import base64
import io
import re
import shutil
import tempfile
import zipfile
from pathlib import Path

from loguru import logger
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

SLIDE_W = 12_192_000
SLIDE_H = 6_858_000
SCALE_X = SLIDE_W / 1280
SCALE_Y = SLIDE_H / 720
SVG_NS = "http://www.w3.org/2000/svg"
XLINK_NS = "http://www.w3.org/1999/xlink"
NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
NS_RELS = "http://schemas.openxmlformats.org/package/2006/relationships"
IMAGE_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"


def assemble_pptx(svg_paths: list[Path], output_path: Path, native: bool = False,
                   embed: bool = False) -> Path:
    """Create a PPTX from SVG files.

    Args:
        svg_paths: Sorted list of SVG file paths (one per slide)
        output_path: Where to save the PPTX
        native: If True, use legacy svg2pptx library.
        embed: If True, embed SVG+PNG as images (old default).
    """
    if native:
        return _assemble_legacy_native(svg_paths, output_path)
    if embed:
        return _assemble_svg_embed(svg_paths, output_path)
    return _assemble_native_shapes(svg_paths, output_path)


# ── Native Shapes Mode (Default) ──────────────────────


def _assemble_native_shapes(svg_paths: list[Path], output_path: Path) -> Path:
    """Convert SVG elements to native DrawingML shapes (directly editable)."""
    from edupptx.output.svg_to_shapes import convert_svg_to_slide_shapes

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    for _ in svg_paths:
        prs.slides.add_slide(prs.slide_layouts[6])

    tmp_dir = Path(tempfile.mkdtemp())
    try:
        base_pptx = tmp_dir / "base.pptx"
        prs.save(str(base_pptx))

        extract_dir = tmp_dir / "content"
        with zipfile.ZipFile(base_pptx, "r") as zf:
            zf.extractall(extract_dir)

        media_dir = extract_dir / "ppt" / "media"
        media_dir.mkdir(exist_ok=True)
        rels_dir = extract_dir / "ppt" / "slides" / "_rels"
        rels_dir.mkdir(exist_ok=True)

        any_images = False

        for i, svg_path in enumerate(svg_paths, 1):
            try:
                slide_xml, media_files, rel_entries = convert_svg_to_slide_shapes(
                    svg_path, slide_num=i
                )

                # Write slide XML
                (extract_dir / "ppt" / "slides" / f"slide{i}.xml").write_text(
                    slide_xml, encoding="utf-8"
                )

                # Write media files
                for media_name, media_data in media_files.items():
                    (media_dir / media_name).write_bytes(media_data)
                    any_images = True

                # Build rels XML
                extra_rels = ""
                for rel in rel_entries:
                    extra_rels += (
                        f'\n  <Relationship Id="{rel["id"]}" '
                        f'Type="{rel["type"]}" Target="{rel["target"]}"/>'
                    )

                rels_xml = (
                    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
                    f'<Relationships xmlns="{NS_RELS}">\n'
                    '  <Relationship Id="rId1" '
                    'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" '
                    f'Target="../slideLayouts/slideLayout7.xml"/>{extra_rels}\n'
                    '</Relationships>'
                )
                (rels_dir / f"slide{i}.xml.rels").write_text(rels_xml, encoding="utf-8")

                logger.debug("[{}/{}] {} → native shapes", i, len(svg_paths), svg_path.name)

            except Exception as e:
                logger.warning("Native shapes failed for {}: {} — falling back to SVG embed", svg_path.name, e)
                # Fallback: embed as SVG+PNG image for this slide
                _embed_single_slide_fallback(extract_dir, media_dir, rels_dir, svg_path, i)

        # Update Content_Types
        ct_path = extract_dir / "[Content_Types].xml"
        ct = ct_path.read_text(encoding="utf-8")
        if any_images:
            for ext, ctype in [("png", "image/png"), ("jpg", "image/jpeg"),
                               ("jpeg", "image/jpeg"), ("gif", "image/gif")]:
                if f'Extension="{ext}"' not in ct:
                    ct = ct.replace("</Types>", f'  <Default Extension="{ext}" ContentType="{ctype}"/>\n</Types>')
        ct_path.write_text(ct, encoding="utf-8")

        # Repackage
        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for fp in extract_dir.rglob("*"):
                if fp.is_file():
                    zf.write(fp, fp.relative_to(extract_dir))

        return output_path
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def _embed_single_slide_fallback(extract_dir: Path, media_dir: Path,
                                  rels_dir: Path, svg_path: Path, slide_num: int):
    """Fallback: embed a single slide as SVG+PNG image when native conversion fails."""
    svg_name = f"image{slide_num}.svg"
    png_name = f"image{slide_num}.png"
    shutil.copy(svg_path, media_dir / svg_name)
    png_ok = _svg_to_png(svg_path, media_dir / png_name)

    png_rid = "rId2" if png_ok else ""
    svg_rid = "rId3" if png_ok else "rId2"

    slide_xml = _make_slide_xml(slide_num, png_rid, svg_rid, png_ok)
    (extract_dir / "ppt" / "slides" / f"slide{slide_num}.xml").write_text(slide_xml, encoding="utf-8")

    rels_xml = _make_slide_rels(png_rid, png_name, svg_rid, svg_name, png_ok)
    (rels_dir / f"slide{slide_num}.xml.rels").write_text(rels_xml, encoding="utf-8")

    # Ensure SVG/PNG content types
    ct_path = extract_dir / "[Content_Types].xml"
    ct = ct_path.read_text(encoding="utf-8")
    if 'Extension="svg"' not in ct:
        ct = ct.replace("</Types>", '  <Default Extension="svg" ContentType="image/svg+xml"/>\n</Types>')
    if 'Extension="png"' not in ct:
        ct = ct.replace("</Types>", '  <Default Extension="png" ContentType="image/png"/>\n</Types>')
    ct_path.write_text(ct, encoding="utf-8")


# ── SVG+PNG Embedding (--embed Mode) ────────────────────


def _assemble_svg_embed(svg_paths: list[Path], output_path: Path) -> Path:
    """Embed SVG files into PPTX with PNG fallback via ZIP post-processing.

    Based on the ppt-master approach: create base PPTX, extract ZIP, inject
    SVG+PNG media and update slide XML with asvg:svgBlip extension.
    """
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    for _ in svg_paths:
        prs.slides.add_slide(prs.slide_layouts[6])

    tmp_dir = Path(tempfile.mkdtemp())
    try:
        base_pptx = tmp_dir / "base.pptx"
        prs.save(str(base_pptx))

        extract_dir = tmp_dir / "content"
        with zipfile.ZipFile(base_pptx, "r") as zf:
            zf.extractall(extract_dir)

        media_dir = extract_dir / "ppt" / "media"
        media_dir.mkdir(exist_ok=True)

        for i, svg_path in enumerate(svg_paths, 1):
            svg_name = f"image{i}.svg"
            png_name = f"image{i}.png"
            png_rid = "rId2"
            svg_rid = "rId3"

            # Copy SVG to media
            shutil.copy(svg_path, media_dir / svg_name)

            # Generate PNG fallback
            png_ok = _svg_to_png(svg_path, media_dir / png_name)
            if not png_ok:
                svg_rid = "rId2"
                png_rid = ""

            # Write slide XML
            slide_xml = _make_slide_xml(i, png_rid, svg_rid, png_ok)
            (extract_dir / "ppt" / "slides" / f"slide{i}.xml").write_text(slide_xml, encoding="utf-8")

            # Write slide rels
            rels_dir = extract_dir / "ppt" / "slides" / "_rels"
            rels_dir.mkdir(exist_ok=True)
            rels_xml = _make_slide_rels(png_rid, png_name, svg_rid, svg_name, png_ok)
            (rels_dir / f"slide{i}.xml.rels").write_text(rels_xml, encoding="utf-8")

            logger.debug("[{}/{}] {} (SVG{})", i, len(svg_paths), svg_path.name, "+PNG" if png_ok else "")

        # Update Content_Types
        ct_path = extract_dir / "[Content_Types].xml"
        ct = ct_path.read_text(encoding="utf-8")
        if 'Extension="svg"' not in ct:
            ct = ct.replace("</Types>", '  <Default Extension="svg" ContentType="image/svg+xml"/>\n</Types>')
        if 'Extension="png"' not in ct:
            ct = ct.replace("</Types>", '  <Default Extension="png" ContentType="image/png"/>\n</Types>')
        ct_path.write_text(ct, encoding="utf-8")

        # Repackage
        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for fp in extract_dir.rglob("*"):
                if fp.is_file():
                    zf.write(fp, fp.relative_to(extract_dir))

        return output_path
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def _svg_to_png(svg_path: Path, png_path: Path) -> bool:
    try:
        import cairosvg
        cairosvg.svg2png(url=str(svg_path), write_to=str(png_path),
                         output_width=2560, output_height=1440)
        return True
    except Exception as e:
        logger.warning("SVG→PNG failed for {}: {}", svg_path.name, e)
        return False


def _make_slide_xml(slide_num: int, png_rid: str, svg_rid: str, has_png: bool) -> str:
    if has_png:
        blip = f'''<a:blip r:embed="{png_rid}">
            <a:extLst>
              <a:ext uri="{{96DAC541-7B7A-43D3-8B79-37D633B846F1}}">
                <asvg:svgBlip xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main" r:embed="{svg_rid}"/>
              </a:ext>
            </a:extLst>
          </a:blip>'''
    else:
        blip = f'<a:blip r:embed="{svg_rid}"/>'

    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
       xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree>
    <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
    <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
    <p:pic>
      <p:nvPicPr><p:cNvPr id="2" name="SVG {slide_num}"/><p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>
      <p:blipFill>{blip}<a:stretch><a:fillRect/></a:stretch></p:blipFill>
      <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="{SLIDE_W}" cy="{SLIDE_H}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
    </p:pic>
  </p:spTree></p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>'''


def _make_slide_rels(png_rid: str, png_name: str, svg_rid: str, svg_name: str, has_png: bool) -> str:
    rels = [
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout7.xml"/>',
    ]
    if has_png:
        rels.append(f'<Relationship Id="{png_rid}" Type="{IMAGE_REL_TYPE}" Target="../media/{png_name}"/>')
    rels.append(f'<Relationship Id="{svg_rid}" Type="{IMAGE_REL_TYPE}" Target="../media/{svg_name}"/>')
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="{NS_RELS}">
  {'  '.join(rels)}
</Relationships>'''


# ── Legacy Native Mode (svg2pptx library) ──────────────


def _assemble_legacy_native(svg_paths: list[Path], output_path: Path) -> Path:
    """Convert SVG to native editable shapes via svg2pptx library (legacy)."""
    from svg2pptx import SVGConverter

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    converter = SVGConverter()

    for svg_path in svg_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        try:
            converter.add_to_slide(str(svg_path), slide)
            _embed_images_from_svg(slide, svg_path)
            _fix_textbox_widths(prs, slide)
            logger.debug("Native: {} ({} shapes)", svg_path.name, len(slide.shapes))
        except Exception as e:
            logger.warning("svg2pptx failed for {}: {}", svg_path.name, e)
            _add_png_fallback(slide, svg_path)

    prs.save(str(output_path))
    return output_path


def _embed_images_from_svg(slide, svg_path: Path) -> None:
    """Parse SVG <image> elements and embed as PPTX pictures."""
    try:
        content = svg_path.read_text(encoding="utf-8")
        content = re.sub(r"&(?!amp;|lt;|gt;|quot;|apos;|#)", "&amp;", content)
        root = etree.fromstring(content.encode("utf-8"))
    except Exception:
        return

    for img in root.iter(f"{{{SVG_NS}}}image"):
        href = img.get("href") or img.get(f"{{{XLINK_NS}}}href") or ""
        if not href:
            continue
        x, y = _pf(img.get("x")), _pf(img.get("y"))
        w, h = _pf(img.get("width"), 200), _pf(img.get("height"), 150)

        img_data = None
        if href.startswith("data:image"):
            try:
                _, b64 = href.split(",", 1)
                img_data = base64.b64decode(b64)
            except Exception:
                continue
        elif Path(href).exists():
            img_data = Path(href).read_bytes()

        if img_data:
            try:
                slide.shapes.add_picture(
                    io.BytesIO(img_data),
                    Emu(int(x * SCALE_X)), Emu(int(y * SCALE_Y)),
                    Emu(int(w * SCALE_X)), Emu(int(h * SCALE_Y)),
                )
            except Exception as e:
                logger.debug("Image embed failed: {}", e)


def _fix_textbox_widths(prs: Presentation, slide) -> None:
    """Expand TextBox widths for CJK text."""
    from pptx.util import Pt
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        fs_emu = Pt(16)
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    fs_emu = run.font.size
                    break
        fs_pt = fs_emu / 12700
        lines = text.split("\n")
        max_line = max(lines, key=len)
        needed = int(len(max_line) * fs_pt * 0.75 * 12700)
        if needed > shape.width:
            max_w = prs.slide_width - shape.left - Emu(50000)
            shape.width = min(needed, max(max_w, shape.width))


def _pf(s: str | None, default: float = 0) -> float:
    if not s:
        return default
    try:
        return float(s.replace("px", "").strip())
    except (ValueError, TypeError):
        return default


def _add_png_fallback(slide, svg_path: Path) -> None:
    try:
        import cairosvg
        png = cairosvg.svg2png(url=str(svg_path), output_width=2560, output_height=1440)
        slide.shapes.add_picture(io.BytesIO(png), Emu(0), Emu(0), Emu(SLIDE_W), Emu(SLIDE_H))
    except Exception as e:
        logger.error("PNG fallback failed for {}: {}", svg_path.name, e)
