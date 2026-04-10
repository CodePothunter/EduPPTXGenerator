"""Native pptx diagram rendering — vector shapes directly on slides."""

from __future__ import annotations

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

from edupptx.design_system import DesignTokens
from edupptx.layout_engine import SlotPosition

_NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def _hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_text_shape(slide, text: str, x: int, y: int, w: int, h: int,
                    font_size: int, color: str, bold: bool = False,
                    font_name: str = "Arial", font_ea: str = "Noto Sans SC",
                    align: PP_ALIGN = PP_ALIGN.CENTER):
    """Add a text box with proper CJK font support."""
    txbox = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = _hex_to_rgb(color)
    font.name = font_name
    # East Asian font
    rPr = run._r.get_or_add_rPr()
    ns = _NSMAP["a"]
    for tag in ["ea", "cs", "sym"]:
        el = rPr.find(f"{{{ns}}}{tag}")
        if el is None:
            el = etree.SubElement(rPr, f"{{{ns}}}{tag}")
        el.set("typeface", font_ea)


def _add_rounded_box(slide, x: int, y: int, w: int, h: int,
                     fill_color: str, border_color: str, radius: int = 5000):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        Emu(x), Emu(y), Emu(w), Emu(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    shape.line.color.rgb = _hex_to_rgb(border_color)
    shape.line.width = Pt(1.5)
    # Adjust corner radius
    ns = _NSMAP["a"]
    prst = shape._element.find(f".//{{{ns}}}prstGeom")
    if prst is not None:
        av_lst = prst.find(f"{{{ns}}}avLst")
        if av_lst is None:
            av_lst = etree.SubElement(prst, f"{{{ns}}}avLst")
        for child in list(av_lst):
            av_lst.remove(child)
        gd = etree.SubElement(av_lst, f"{{{ns}}}gd")
        gd.set("name", "adj")
        gd.set("fmla", f"val {radius}")
    return shape


def _add_line(slide, x1: int, y1: int, x2: int, y2: int, color: str, width_pt: float = 2):
    """Add a line connector."""
    from pptx.util import Emu, Pt
    left = min(x1, x2)
    top = min(y1, y2)
    w = abs(x2 - x1) or 12700  # min 1pt width for vertical lines
    h = abs(y2 - y1) or 12700
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Emu(x1), Emu(y1), Emu(x2), Emu(y2),
    )
    connector.line.color.rgb = _hex_to_rgb(color)
    connector.line.width = Pt(width_pt)


def _add_circle(slide, cx: int, cy: int, radius: int, fill_color: str):
    """Add a filled circle."""
    shape = slide.shapes.add_shape(
        9,  # MSO_SHAPE.OVAL
        Emu(cx - radius), Emu(cy - radius), Emu(radius * 2), Emu(radius * 2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    shape.line.fill.background()
    return shape


# ── Diagram Renderers ──────────────────────────────────────

def _draw_flowchart(slide, data: dict, slot: SlotPosition, design: DesignTokens):
    """Draw flowchart with native shapes: rounded boxes + line connectors."""
    nodes = data.get("nodes", [])
    edges = data.get("edges", [])
    direction = data.get("direction", "TB")
    if not nodes:
        return

    n = len(nodes)
    box_w = 2032000   # 160pt
    box_h = 508000    # 40pt
    gap = 381000      # 30pt

    if direction == "TB":
        total_h = n * box_h + (n - 1) * gap
        start_y = slot.y + (slot.height - total_h) // 2
        cx = slot.x + slot.width // 2

        node_positions = {}
        for i, node in enumerate(nodes):
            x = cx - box_w // 2
            y = start_y + i * (box_h + gap)
            node_positions[node.get("id", str(i))] = (x, y)
            _add_rounded_box(slide, x, y, box_w, box_h,
                             design.accent_light, design.accent)
            _add_text_shape(slide, node["label"],
                            x, y, box_w, box_h,
                            font_size=12, color=design.text_primary, bold=True)

        for edge in edges:
            src = node_positions.get(edge.get("from"))
            dst = node_positions.get(edge.get("to"))
            if src and dst:
                _add_line(slide,
                          src[0] + box_w // 2, src[1] + box_h,
                          dst[0] + box_w // 2, dst[1],
                          design.accent, width_pt=2)
    else:  # LR
        total_w = n * box_w + (n - 1) * gap
        start_x = slot.x + (slot.width - total_w) // 2
        cy = slot.y + slot.height // 2

        node_positions = {}
        for i, node in enumerate(nodes):
            x = start_x + i * (box_w + gap)
            y = cy - box_h // 2
            node_positions[node.get("id", str(i))] = (x, y)
            _add_rounded_box(slide, x, y, box_w, box_h,
                             design.accent_light, design.accent)
            _add_text_shape(slide, node["label"],
                            x, y, box_w, box_h,
                            font_size=12, color=design.text_primary, bold=True)

        for edge in edges:
            src = node_positions.get(edge.get("from"))
            dst = node_positions.get(edge.get("to"))
            if src and dst:
                _add_line(slide,
                          src[0] + box_w, src[1] + box_h // 2,
                          dst[0], dst[1] + box_h // 2,
                          design.accent, width_pt=2)


def _draw_timeline(slide, data: dict, slot: SlotPosition, design: DesignTokens):
    """Draw timeline with native shapes: circles on a line + labels."""
    events = data.get("events", [])
    if not events:
        return

    n = len(events)
    padding = 508000  # 40pt side padding
    line_y = slot.y + slot.height * 2 // 5
    line_x1 = slot.x + padding
    line_x2 = slot.x + slot.width - padding
    line_w = line_x2 - line_x1

    # Horizontal line
    _add_line(slide, line_x1, line_y, line_x2, line_y, design.accent, width_pt=3)

    # Events
    dot_r = 76200  # 6pt radius
    label_w = line_w // n
    for i, event in enumerate(events):
        cx = line_x1 + (line_w * i) // max(n - 1, 1) if n > 1 else line_x1 + line_w // 2
        # Dot
        _add_circle(slide, cx, line_y, dot_r, design.accent)
        # Year label above
        year = event.get("year", "")
        _add_text_shape(slide, year,
                        cx - label_w // 2, line_y - 381000,  # 30pt above
                        label_w, 254000,  # 20pt high
                        font_size=10, color=design.accent, bold=True)
        # Description below
        label = event.get("label", "")
        desc = event.get("description", "")
        text = label if not desc else f"{label}\n{desc}"
        _add_text_shape(slide, text,
                        cx - label_w // 2, line_y + 127000,  # 10pt below
                        label_w, 508000,  # 40pt high
                        font_size=9, color=design.text_secondary)


def _draw_comparison(slide, data: dict, slot: SlotPosition, design: DesignTokens):
    """Draw comparison columns with native shapes."""
    columns = data.get("columns", [])
    if not columns:
        return

    n = len(columns)
    col_gap = 127000  # 10pt
    col_w = (slot.width - (n - 1) * col_gap) // n
    header_h = 508000  # 40pt

    for i, col in enumerate(columns):
        x = slot.x + i * (col_w + col_gap)
        # Header box
        _add_rounded_box(slide, x, slot.y, col_w, header_h,
                         design.accent, design.accent, radius=3000)
        _add_text_shape(slide, col.get("header", ""),
                        x, slot.y, col_w, header_h,
                        font_size=13, color="#FFFFFF", bold=True)
        # Items
        items = col.get("items", [])
        item_h = 254000  # 20pt per item
        for j, item in enumerate(items):
            iy = slot.y + header_h + 127000 + j * (item_h + 63500)
            _add_text_shape(slide, f"  {item}",
                            x + 63500, iy, col_w - 127000, item_h,
                            font_size=11, color=design.text_primary,
                            align=PP_ALIGN.LEFT)


def _draw_hierarchy(slide, data: dict, slot: SlotPosition, design: DesignTokens):
    """Draw hierarchy tree with native shapes."""
    root = data.get("root")
    if not root:
        return

    def _count_leaves(node):
        children = node.get("children", [])
        if not children:
            return 1
        return sum(_count_leaves(c) for c in children)

    def _draw_node(node, x, y, available_w, depth):
        box_w = min(1524000, available_w - 254000)  # max 120pt, fit in space
        box_h = 381000  # 30pt
        bx = x + (available_w - box_w) // 2

        _add_rounded_box(slide, bx, y, box_w, box_h,
                         design.accent_light if depth > 0 else design.accent,
                         design.accent, radius=4000)
        _add_text_shape(slide, node.get("label", ""),
                        bx, y, box_w, box_h,
                        font_size=11,
                        color="#FFFFFF" if depth == 0 else design.text_primary,
                        bold=True)

        children = node.get("children", [])
        if not children:
            return

        child_y = y + box_h + 254000  # 20pt gap
        total_leaves = _count_leaves(node)
        child_x = x
        parent_cx = bx + box_w // 2

        for child in children:
            leaves = _count_leaves(child)
            child_w = (available_w * leaves) // total_leaves
            child_cx = child_x + child_w // 2

            # Line from parent bottom to child top
            _add_line(slide, parent_cx, y + box_h, child_cx, child_y,
                      design.accent, width_pt=1.5)

            _draw_node(child, child_x, child_y, child_w, depth + 1)
            child_x += child_w

    _draw_node(root, slot.x, slot.y, slot.width, 0)


def _draw_cycle(slide, data: dict, slot: SlotPosition, design: DesignTokens):
    """Draw cycle diagram with native shapes arranged in a circle."""
    import math
    steps = data.get("steps", [])
    if not steps:
        return

    n = len(steps)
    cx = slot.x + slot.width // 2
    cy = slot.y + slot.height // 2
    radius = min(slot.width, slot.height) * 2 // 5
    box_w = 1524000  # 120pt
    box_h = 381000   # 30pt

    positions = []
    for i in range(n):
        angle = -math.pi / 2 + 2 * math.pi * i / n
        px = int(cx + radius * math.cos(angle))
        py = int(cy + radius * math.sin(angle))
        positions.append((px, py))

        # Box
        bx = px - box_w // 2
        by = py - box_h // 2
        _add_rounded_box(slide, bx, by, box_w, box_h,
                         design.accent_light, design.accent)
        _add_text_shape(slide, steps[i].get("label", ""),
                        bx, by, box_w, box_h,
                        font_size=11, color=design.text_primary, bold=True)

    # Arrows between steps (lines from edge of one box toward next)
    for i in range(n):
        j = (i + 1) % n
        x1, y1 = positions[i]
        x2, y2 = positions[j]
        # Shorten line to not overlap boxes
        dx, dy = x2 - x1, y2 - y1
        dist = math.sqrt(dx * dx + dy * dy)
        if dist == 0:
            continue
        ux, uy = dx / dist, dy / dist
        offset = 254000  # 20pt
        lx1 = int(x1 + ux * offset)
        ly1 = int(y1 + uy * offset)
        lx2 = int(x2 - ux * offset)
        ly2 = int(y2 - uy * offset)
        _add_line(slide, lx1, ly1, lx2, ly2, design.accent, width_pt=1.5)


# ── Dispatcher ──────────────────────────────────────

_RENDERERS = {
    "flowchart": _draw_flowchart,
    "timeline": _draw_timeline,
    "comparison": _draw_comparison,
    "hierarchy": _draw_hierarchy,
    "cycle": _draw_cycle,
}


def draw_diagram_on_slide(
    slide, diagram_type: str, data: dict,
    slot: SlotPosition, design: DesignTokens,
) -> None:
    """Draw a diagram using native pptx shapes directly on the slide."""
    renderer = _RENDERERS.get(diagram_type)
    if renderer:
        renderer(slide, data, slot, design)
