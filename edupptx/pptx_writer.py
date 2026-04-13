"""PPTX writer: dumb shape placer for the v2 pipeline.

Takes list[ResolvedSlide] and writes a .pptx file. No style decisions —
all values come from ResolvedShape fields. Each _write_* method is 15-30 lines.
"""

from __future__ import annotations

import io
import os
import tempfile
from pathlib import Path

from loguru import logger
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from edupptx.diagram_native import SlotPosition, draw_diagram_on_slide
from edupptx.icons import get_icon_png, get_icon_svg
from edupptx.models import ResolvedShape, ResolvedSlide
from edupptx.style_schema import SLIDE_H, SLIDE_W, ResolvedStyle
from edupptx.xml_patches import (
    hex_to_rgb,
    patch_alpha,
    patch_autofit,
    patch_corner_radius,
    patch_shadow,
    patch_v_anchor,
    set_font_cjk,
    NSMAP,
)

from lxml import etree


class PptxWriter:
    """Writes ResolvedSlide objects to a .pptx file."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Emu(SLIDE_W)
        self.prs.slide_height = Emu(SLIDE_H)
        self._temp_dir = tempfile.mkdtemp(prefix="edupptx_v2_")

    def write_slides(self, slides: list[ResolvedSlide],
                     bg_paths: list[Path] | None = None,
                     style: ResolvedStyle | None = None) -> None:
        """Write all resolved slides into the presentation."""
        for i, rs in enumerate(slides):
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank

            # Background
            bg = rs.background_path
            if bg is None and bg_paths:
                bg = bg_paths[i % len(bg_paths)]
            if bg and bg.exists():
                slide.shapes.add_picture(
                    str(bg), Emu(0), Emu(0), Emu(SLIDE_W), Emu(SLIDE_H),
                )

            # Shapes sorted by z_order
            for shape in sorted(rs.shapes, key=lambda s: s.z_order):
                self._write_shape(slide, shape)

            # Native diagram (vector shapes drawn directly on slide)
            if rs.diagram_info and style:
                d_type, d_data, (sx, sy, sw, sh) = rs.diagram_info
                draw_diagram_on_slide(
                    slide, d_type, d_data,
                    SlotPosition(sx, sy, sw, sh), style,
                )

            # Speaker notes
            if rs.notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = rs.notes

    def save(self, path: str | Path) -> Path:
        out = Path(path)
        self.prs.save(str(out))
        logger.info("V2 pipeline saved: {}", out)
        return out

    def _write_shape(self, slide, shape: ResolvedShape) -> None:
        match shape.shape_type:
            case "textbox":
                self._write_textbox(slide, shape)
            case "rounded_rect":
                self._write_rounded_rect(slide, shape)
            case "oval":
                self._write_oval(slide, shape)
            case "image":
                self._write_image(slide, shape)
            case "line":
                self._write_line(slide, shape)
            case _:
                logger.warning("Unknown shape type '{}', skipping", shape.shape_type)

    def _write_textbox(self, slide, shape: ResolvedShape) -> None:
        txbox = slide.shapes.add_textbox(
            Emu(shape.left), Emu(shape.top),
            Emu(shape.width), Emu(shape.height),
        )
        tf = txbox.text_frame
        tf.word_wrap = True

        # Vertical anchor + autofit
        patch_v_anchor(txbox, shape.v_anchor)
        if shape.auto_shrink:
            patch_autofit(txbox)

        p = tf.paragraphs[0]
        # Determine alignment from v_anchor or font bold as heuristic
        if shape.font and shape.font.bold and shape.v_anchor == "ctr":
            p.alignment = PP_ALIGN.CENTER
        elif shape.v_anchor == "ctr":
            p.alignment = PP_ALIGN.CENTER
        else:
            p.alignment = PP_ALIGN.LEFT

        run = p.add_run()
        run.text = shape.text or ""
        if shape.font:
            set_font_cjk(
                run, shape.font.family, shape.font.fallback,
                shape.font.size_pt, shape.font.bold, shape.font.color,
            )

    def _write_rounded_rect(self, slide, shape: ResolvedShape) -> None:
        s = slide.shapes.add_shape(
            5,  # MSO_SHAPE.ROUNDED_RECTANGLE
            Emu(shape.left), Emu(shape.top),
            Emu(shape.width), Emu(shape.height),
        )

        if shape.fill_color:
            s.fill.solid()
            s.fill.fore_color.rgb = hex_to_rgb(shape.fill_color)
        else:
            s.fill.background()

        if shape.line_color:
            s.line.color.rgb = hex_to_rgb(shape.line_color)
            s.line.width = Pt(1.5)
        else:
            s.line.fill.background()

        if shape.corner_radius:
            patch_corner_radius(s, shape.corner_radius)

        if shape.shadow:
            patch_shadow(s, shape.shadow)

        if shape.alpha_pct < 100:
            patch_alpha(s, shape.alpha_pct)

    def _write_oval(self, slide, shape: ResolvedShape) -> None:
        s = slide.shapes.add_shape(
            9,  # MSO_SHAPE.OVAL
            Emu(shape.left), Emu(shape.top),
            Emu(shape.width), Emu(shape.height),
        )
        if shape.fill_color:
            s.fill.solid()
            s.fill.fore_color.rgb = hex_to_rgb(shape.fill_color)
        else:
            s.fill.background()
        s.line.fill.background()

    def _write_image(self, slide, shape: ResolvedShape) -> None:
        if not shape.image_path:
            return

        # Icon images use "icon:name:color" format
        if shape.image_path.startswith("icon:"):
            parts = shape.image_path.split(":", 2)
            icon_name = parts[1]
            icon_color = parts[2] if len(parts) > 2 else "#000000"
            self._write_icon(slide, icon_name, icon_color, shape)
        elif Path(shape.image_path).exists():
            slide.shapes.add_picture(
                shape.image_path,
                Emu(shape.left), Emu(shape.top),
                Emu(shape.width), Emu(shape.height),
            )

    def _write_icon(self, slide, icon_name: str, icon_color: str,
                    shape: ResolvedShape) -> None:
        """Write an icon image with optional SVG blip extension."""

        png_bytes = get_icon_png(icon_name, icon_color, size=48)
        png_stream = io.BytesIO(png_bytes)

        pic = slide.shapes.add_picture(
            png_stream,
            Emu(shape.left), Emu(shape.top),
            Emu(shape.width), Emu(shape.height),
        )

        # Try SVG extension for modern PowerPoint
        try:
            svg_str = get_icon_svg(icon_name, icon_color)
            svg_path = os.path.join(self._temp_dir, f"{icon_name}.svg")
            with open(svg_path, "wb") as f:
                f.write(svg_str.encode("utf-8"))

            slide_part = pic.part
            svg_rel = slide_part.relate_to_file(
                svg_path,
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
            )

            ns_a = NSMAP["a"]
            blip = pic._element.find(f".//{{{ns_a}}}blip")
            if blip is not None:
                ext_lst = blip.find(f"{{{ns_a}}}extLst")
                if ext_lst is None:
                    ext_lst = etree.SubElement(blip, f"{{{ns_a}}}extLst")
                ext = etree.SubElement(ext_lst, f"{{{ns_a}}}ext")
                ext.set("uri", "{96DAC541-7B7A-43D3-8B79-37D633B846F1}")
                ns_asvg = NSMAP["asvg"]
                svg_blip = etree.SubElement(ext, f"{{{ns_asvg}}}svgBlip")
                svg_blip.set(f"{{{NSMAP['r']}}}embed", svg_rel)
        except Exception:
            pass  # PNG fallback

    def _write_line(self, slide, shape: ResolvedShape) -> None:
        """Write a simple horizontal line."""
        s = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Emu(shape.left), Emu(shape.top),
            Emu(shape.width), Emu(shape.height),
        )
        if shape.fill_color:
            s.fill.solid()
            s.fill.fore_color.rgb = hex_to_rgb(shape.fill_color)
        s.line.fill.background()
