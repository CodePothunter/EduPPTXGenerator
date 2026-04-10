"""OOXML renderer — python-pptx + XML patches for SVG icons, shadows, alpha."""

from __future__ import annotations

import io
import os
import tempfile
from copy import deepcopy
from pathlib import Path
from typing import BinaryIO

from loguru import logger
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Emu, Pt

from edupptx.design_system import DesignTokens
from edupptx.icons import get_icon_png, get_icon_svg
from edupptx.layout_engine import (
    CARD_GAP,
    CARD_TOP,
    FOOTER_H,
    FOOTER_Y,
    MARGIN_X,
    MARGIN_Y,
    CONTENT_W,
    SLIDE_H,
    SLIDE_W,
    TITLE_Y,
    SlotLayout,
    SlotPosition,
    get_layout,
)
from edupptx.models import PresentationPlan, SlideCard, SlideContent

# Slide types that get special visual treatment (no standard title underline)
_NO_UNDERLINE_TYPES = {"cover", "closing", "section", "big_quote"}

# Slide types that skip the content panel (use background directly)
_NO_PANEL_TYPES = {"closing", "section", "big_quote", "full_image"}

# XML namespaces
_NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a14": "http://schemas.microsoft.com/office/drawing/2010/main",
    "asvg": "http://schemas.microsoft.com/office/drawing/2016/SVG/main",
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_font(run, design: DesignTokens, size_pt: int, bold: bool = False, color: str | None = None):
    """Configure font properties on a text run."""
    font = run.font
    font.size = Pt(size_pt)
    font.bold = bold
    font.color.rgb = _hex_to_rgb(color or design.text_primary)
    font.name = design.font_fallback  # Latin typeface

    # Set East Asian / Complex Script / Symbol typefaces via XML
    rPr = run._r.get_or_add_rPr()
    for tag, attr in [("a:ea", design.font_primary), ("a:cs", design.font_primary), ("a:sym", design.font_primary)]:
        ns = _NSMAP["a"]
        el = rPr.find(f"{{{ns}}}{tag.split(':')[1]}")
        if el is None:
            el = etree.SubElement(rPr, f"{{{ns}}}{tag.split(':')[1]}")
        el.set("typeface", attr)


class PresentationRenderer:
    """Renders a PresentationPlan into a .pptx file."""

    def __init__(self, design: DesignTokens):
        self.prs = Presentation()
        self.prs.slide_width = Emu(SLIDE_W)
        self.prs.slide_height = Emu(SLIDE_H)
        self.design = design
        self._temp_dir = tempfile.mkdtemp(prefix="edupptx_")

    def render(self, plan: PresentationPlan, backgrounds: list[Path]) -> Path:
        """Render all slides and save to a .pptx file."""
        for i, slide_content in enumerate(plan.slides):
            bg = backgrounds[i % len(backgrounds)]
            logger.info("Rendering slide %d/%d: %s", i + 1, len(plan.slides), slide_content.type)
            self._render_slide(slide_content, bg)

        output = Path(f"{plan.topic}.pptx")
        self.prs.save(str(output))
        logger.info("Saved presentation: %s", output)
        return output

    def save(self, path: str | Path) -> Path:
        out = Path(path)
        self.prs.save(str(out))
        return out

    def render_slide(
        self, content: SlideContent,
        bg_path: Path | None = None,
        material_path: Path | None = None,
    ) -> None:
        """Render a single slide into the presentation."""
        if bg_path is None:
            from PIL import Image
            import tempfile
            img = Image.new("RGB", (1920, 1080), tuple(int(self.design.bg_overlay.lstrip('#')[i:i+2], 16) for i in (0, 2, 4)))
            bg_path = Path(tempfile.mktemp(suffix=".jpg"))
            img.save(bg_path, "JPEG")
        self._render_slide(content, bg_path, material_path)

    def _render_slide(self, content: SlideContent, bg_path: Path, material_path: Path | None = None):
        """Render a single slide with all its components."""
        layout_idx = min(6, len(self.prs.slide_layouts) - 1)
        layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(layout)

        # Determine material position from content_materials
        mat_position = None
        if material_path and material_path.exists() and content.content_materials:
            mat_position = content.content_materials[0].position

        slots = get_layout(content.type, len(content.cards), material_position=mat_position)

        # 1. Background image (full canvas)
        self._add_background(slide, bg_path, slots.background)

        # 1.5. Content area panel (frosted glass effect for depth)
        if content.type not in _NO_PANEL_TYPES and content.cards:
            self._add_content_panel(slide, slots)

        # 2. Special page decorations (behind content)
        if content.type == "big_quote":
            self._add_quote_decoration(slide, slots.title)
        elif content.type == "section":
            self._add_section_decoration(slide, slots.title)
        elif content.type == "closing":
            self._add_closing_decoration(slide)

        # 3. Title
        if content.type == "big_quote":
            # Big quote gets special large italic treatment
            self._add_textbox(
                slide, content.title, slots.title,
                size_pt=self.design.size_title + 4, bold=False,
                color=self.design.text_primary,
                align=PP_ALIGN.CENTER,
            )
        elif content.type in ("section", "closing"):
            self._add_textbox(
                slide, content.title, slots.title,
                size_pt=self.design.size_title, bold=True,
                align=PP_ALIGN.CENTER,
            )
        else:
            self._add_textbox(
                slide, content.title, slots.title,
                size_pt=self.design.size_title, bold=True,
            )

        # 4. Title accent underline (content slides only)
        if content.type not in _NO_UNDERLINE_TYPES:
            self._add_title_underline(slide, slots.title)

        # 5. Subtitle
        if content.subtitle and slots.subtitle:
            sub_align = PP_ALIGN.CENTER if content.type in ("section", "closing", "cover") else PP_ALIGN.LEFT
            self._add_textbox(
                slide, content.subtitle, slots.subtitle,
                size_pt=self.design.size_subtitle,
                color=self.design.text_secondary,
                align=sub_align,
            )

        # 6. Content material (diagram or illustration)
        if content.content_materials and slots.material_slot:
            mat = content.content_materials[0]
            if mat.diagram_type and mat.diagram_data:
                from edupptx.diagram_native import draw_diagram_on_slide
                draw_diagram_on_slide(
                    slide, mat.diagram_type, mat.diagram_data,
                    slots.material_slot, self.design,
                )
            elif mat.action == "generate_illustration" and material_path and material_path.exists():
                self._add_illustration(
                    slide, material_path, slots.material_slot,
                    anchor=mat.image_anchor, scale=mat.image_scale,
                )

        # 7. Cards with icons
        for j, card in enumerate(content.cards):
            if j >= len(slots.cards):
                break
            self._add_card(
                slide, card,
                slots.cards[j],
                slots.card_icons[j] if j < len(slots.card_icons) else None,
                slots.card_titles[j] if j < len(slots.card_titles) else None,
                slots.card_bodies[j] if j < len(slots.card_bodies) else None,
            )

        # 8. Formula
        if content.formula and slots.formula:
            self._add_formula_bar(slide, content.formula, slots.formula)

        # 9. Footer
        if content.footer and slots.footer:
            self._add_footer(slide, content.footer, slots.footer)

        # 10. Speaker notes
        if content.notes:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = content.notes

    def _add_background(self, slide, bg_path: Path, slot: SlotPosition):
        """Add a full-canvas background image."""
        slide.shapes.add_picture(
            str(bg_path),
            Emu(slot.x), Emu(slot.y), Emu(slot.width), Emu(slot.height),
        )

    def _add_overlay(self, slide, slot: SlotPosition):
        """Add a semi-transparent color overlay."""
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Emu(slot.x), Emu(slot.y), Emu(slot.width), Emu(slot.height),
        )
        shape.line.fill.background()  # No border

        # Set fill with alpha via XML
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(self.design.bg_overlay)

        # Patch alpha
        sp_xml = shape._element
        ns_a = _NSMAP["a"]
        solid_fill = sp_xml.find(f".//{{{ns_a}}}solidFill")
        if solid_fill is not None:
            color_el = solid_fill[0] if len(solid_fill) else None
            if color_el is not None:
                alpha_val = int(self.design.bg_overlay_alpha * 100000)
                alpha_el = etree.SubElement(color_el, f"{{{ns_a}}}alpha")
                alpha_el.set("val", str(alpha_val))

    def _add_textbox(
        self, slide, text: str, slot: SlotPosition,
        size_pt: int = 16, bold: bool = False, color: str | None = None,
        align: PP_ALIGN = PP_ALIGN.LEFT,
    ):
        """Add a text box at the given slot position."""
        txbox = slide.shapes.add_textbox(
            Emu(slot.x), Emu(slot.y), Emu(slot.width), Emu(slot.height),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        _set_font(run, self.design, size_pt, bold=bold, color=color)

    def _tint_color(self, hex_color: str, amount: float = 0.12) -> str:
        """Mix hex_color toward white. amount=0 → white, amount=1 → full color."""
        r = int(hex_color.lstrip('#')[0:2], 16)
        g = int(hex_color.lstrip('#')[2:4], 16)
        b = int(hex_color.lstrip('#')[4:6], 16)
        r = int(255 + (r - 255) * amount)
        g = int(255 + (g - 255) * amount)
        b = int(255 + (b - 255) * amount)
        return f"#{r:02X}{g:02X}{b:02X}"

    def _add_card(
        self, slide, card: SlideCard,
        card_slot: SlotPosition,
        icon_slot: SlotPosition | None,
        title_slot: SlotPosition | None,
        body_slot: SlotPosition | None,
    ):
        """Add a card component: container + accent header + icon + title + body."""
        # Card container (rounded rectangle with shadow)
        card_fill = self._tint_color(self.design.accent_light, 0.15)
        shape = slide.shapes.add_shape(
            5,  # MSO_SHAPE.ROUNDED_RECTANGLE
            Emu(card_slot.x), Emu(card_slot.y),
            Emu(card_slot.width), Emu(card_slot.height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(card_fill)
        shape.line.fill.background()

        # Patch: add shadow
        self._patch_card_shadow(shape)

        # Patch: adjust corner radius
        self._patch_corner_radius(shape, 8000)

        # Accent header strip inside card top
        strip_inset = int(card_slot.width * 0.08)  # match corner radius
        strip_h = 50800  # 4pt
        strip_y = card_slot.y + strip_inset
        strip = slide.shapes.add_shape(
            5,  # ROUNDED_RECTANGLE
            Emu(card_slot.x + strip_inset), Emu(strip_y),
            Emu(card_slot.width - 2 * strip_inset), Emu(strip_h),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = _hex_to_rgb(self.design.accent)
        strip.line.fill.background()
        self._patch_corner_radius(strip, 50000)

        # Icon
        if icon_slot:
            self._add_icon(slide, card.icon, icon_slot)

        # Card title — use accent color for hierarchy
        if title_slot:
            self._add_textbox(
                slide, card.title, title_slot,
                size_pt=self.design.size_card_title, bold=True,
                color=self.design.accent,
                align=PP_ALIGN.CENTER,
            )

        # Card body
        if body_slot:
            body_size = self.design.size_card_body
            if len(card.body) > 50:
                body_size = max(8, body_size - 2)
            self._add_textbox(
                slide, card.body, body_slot,
                size_pt=body_size,
                color=self.design.text_secondary,
            )

    def _add_icon(self, slide, icon_name: str, slot: SlotPosition):
        """Add an icon with a circular accent background and outer ring."""
        bg_pad = int(slot.width * 0.20)

        # Outer ring (subtle border effect for depth)
        ring_pad = bg_pad + int(slot.width * 0.06)
        ring = slide.shapes.add_shape(
            9,  # MSO_SHAPE.OVAL
            Emu(slot.x - ring_pad), Emu(slot.y - ring_pad),
            Emu(slot.width + ring_pad * 2), Emu(slot.height + ring_pad * 2),
        )
        ring.fill.solid()
        ring_color = self._tint_color(self.design.accent_light, 0.5)
        ring.fill.fore_color.rgb = _hex_to_rgb(ring_color)
        ring.line.fill.background()

        # Inner circle background
        bg_shape = slide.shapes.add_shape(
            9,  # MSO_SHAPE.OVAL
            Emu(slot.x - bg_pad), Emu(slot.y - bg_pad),
            Emu(slot.width + bg_pad * 2), Emu(slot.height + bg_pad * 2),
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = _hex_to_rgb(self.design.accent_light)
        bg_shape.line.fill.background()

        # Generate PNG
        png_bytes = get_icon_png(icon_name, self.design.icon_color, size=48)
        png_stream = io.BytesIO(png_bytes)

        pic = slide.shapes.add_picture(
            png_stream,
            Emu(slot.x), Emu(slot.y), Emu(slot.width), Emu(slot.height),
        )

        # Try to add SVG extension for modern PowerPoint
        try:
            self._patch_svg_blip(pic, icon_name)
        except Exception:
            pass  # PNG fallback is fine

    def _add_formula_bar(self, slide, formula: str, slot: SlotPosition):
        """Add a formula highlight bar with accent border."""
        # Background shape
        shape = slide.shapes.add_shape(
            5,  # ROUNDED_RECTANGLE
            Emu(slot.x), Emu(slot.y), Emu(slot.width), Emu(slot.height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(self.design.accent_light)
        shape.line.color.rgb = _hex_to_rgb(self.design.accent)
        shape.line.width = Pt(1.5)
        self._patch_corner_radius(shape, 5000)

        # Formula text
        self._add_textbox(
            slide, formula, slot,
            size_pt=self.design.size_formula, bold=True,
            color=self.design.accent,
            align=PP_ALIGN.CENTER,
        )

    def _add_footer(self, slide, text: str, slot: SlotPosition):
        """Add a footer text line with subtle accent separator."""
        # Thin accent line above footer
        line_w = int(CONTENT_W * 0.3)
        line_h = 19050  # 1.5pt
        line_x = MARGIN_X + (CONTENT_W - line_w) // 2
        line_y = slot.y - 76200  # 6pt above footer

        sep = slide.shapes.add_shape(
            5, Emu(line_x), Emu(line_y), Emu(line_w), Emu(line_h),
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = _hex_to_rgb(self.design.accent_light)
        sep.line.fill.background()
        self._patch_corner_radius(sep, 50000)

        self._add_textbox(
            slide, text, slot,
            size_pt=self.design.size_footer,
            color=self.design.text_secondary,
            align=PP_ALIGN.CENTER,
        )

    def _add_title_underline(self, slide, title_slot: SlotPosition):
        """Add a short accent-colored line under the title."""
        line_w = 762000   # 60pt
        line_h = 38100    # 3pt
        line_x = title_slot.x
        line_y = title_slot.y + title_slot.height + 50800  # 4pt below title

        shape = slide.shapes.add_shape(
            5, Emu(line_x), Emu(line_y), Emu(line_w), Emu(line_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(self.design.accent)
        shape.line.fill.background()
        self._patch_corner_radius(shape, 50000)

    def _add_quote_decoration(self, slide, title_slot: SlotPosition):
        """Add large decorative quotation marks for big_quote slides."""
        quote_size = 1524000  # 120pt — very large decorative mark
        quote_x = title_slot.x - 127000  # 10pt left of title
        quote_y = title_slot.y - int(quote_size * 0.6)

        # Large opening quotation mark in accent_light (watermark style)
        self._add_textbox(
            slide, "\u201C",
            SlotPosition(quote_x, quote_y, quote_size, quote_size),
            size_pt=96, bold=True,
            color=self.design.accent_light,
        )

        # Vertical accent bar on the left
        bar_w = 50800    # 4pt
        bar_h = title_slot.height + 254000  # extends beyond title
        bar_x = title_slot.x - 254000  # 20pt left of title
        bar_y = title_slot.y - 127000

        bar = slide.shapes.add_shape(
            5, Emu(bar_x), Emu(bar_y), Emu(bar_w), Emu(bar_h),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _hex_to_rgb(self.design.accent)
        bar.line.fill.background()
        self._patch_corner_radius(bar, 50000)

    def _add_section_decoration(self, slide, title_slot: SlotPosition):
        """Add decorative elements for section transition slides."""
        # Horizontal accent line above the title
        line_w = 1270000  # 100pt
        line_h = 50800    # 4pt
        line_x = (SLIDE_W - line_w) // 2
        line_y = title_slot.y - 254000  # 20pt above title

        line = slide.shapes.add_shape(
            5, Emu(line_x), Emu(line_y), Emu(line_w), Emu(line_h),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _hex_to_rgb(self.design.accent)
        line.line.fill.background()
        self._patch_corner_radius(line, 50000)

        # Small decorative diamond in center
        diamond_size = 127000  # 10pt
        diamond_x = (SLIDE_W - diamond_size) // 2
        diamond_y = line_y - diamond_size - 63500  # 5pt above line

        diamond = slide.shapes.add_shape(
            4,  # MSO_SHAPE.DIAMOND
            Emu(diamond_x), Emu(diamond_y), Emu(diamond_size), Emu(diamond_size),
        )
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = _hex_to_rgb(self.design.accent)
        diamond.line.fill.background()

    def _add_closing_decoration(self, slide):
        """Add decorative elements for closing slides."""
        # Centered accent circle (decorative background element)
        circle_size = 2540000  # 200pt
        circle_x = (SLIDE_W - circle_size) // 2
        circle_y = (SLIDE_H - circle_size) // 2 - 381000  # slightly above center

        circle = slide.shapes.add_shape(
            9,  # MSO_SHAPE.OVAL
            Emu(circle_x), Emu(circle_y), Emu(circle_size), Emu(circle_size),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = _hex_to_rgb(self.design.accent_light)
        circle.line.fill.background()

        # Patch alpha for translucency
        ns_a = _NSMAP["a"]
        solid_fill = circle._element.find(f".//{{{ns_a}}}solidFill")
        if solid_fill is not None and len(solid_fill):
            alpha_el = etree.SubElement(solid_fill[0], f"{{{ns_a}}}alpha")
            alpha_el.set("val", "30000")  # 30% opacity

    def _add_content_panel(self, slide, slots: SlotLayout):
        """Add a semi-transparent rounded panel behind the content area.

        Creates visual depth: background → panel → cards.
        """
        pad = 152400  # 12pt padding around content
        panel_x = MARGIN_X - pad
        panel_y = CARD_TOP - pad
        panel_w = CONTENT_W + pad * 2

        # Panel extends to cover cards + footer area
        if slots.footer:
            panel_bottom = FOOTER_Y + FOOTER_H + pad
        else:
            panel_bottom = CARD_TOP + (slots.cards[0].height if slots.cards else 2540000) + pad
        panel_h = panel_bottom - panel_y

        shape = slide.shapes.add_shape(
            5,  # ROUNDED_RECTANGLE
            Emu(panel_x), Emu(panel_y), Emu(panel_w), Emu(panel_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.fill.background()
        self._patch_corner_radius(shape, 3000)

        # Patch alpha for translucency (25% white)
        ns_a = _NSMAP["a"]
        solid_fill = shape._element.find(f".//{{{ns_a}}}solidFill")
        if solid_fill is not None and len(solid_fill):
            alpha_el = etree.SubElement(solid_fill[0], f"{{{ns_a}}}alpha")
            alpha_el.set("val", "25000")  # 25% opacity

    def _add_illustration(
        self, slide, img_path: Path, slot: SlotPosition,
        anchor: str = "center", scale: float = 0.95,
    ):
        """Add an illustration image, maintaining aspect ratio.

        anchor: 'top' | 'center' | 'bottom' — vertical alignment within slot.
        scale:  0.4-1.0 — how much of the slot the image fills.
        """
        try:
            from PIL import Image as PILImage
            with PILImage.open(img_path) as img:
                img_w, img_h = img.size

            # Scale the available area
            avail_w = int(slot.width * scale)
            avail_h = int(slot.height * scale)

            # Fit image within the scaled area, preserving aspect ratio
            img_ratio = img_w / img_h
            avail_ratio = avail_w / avail_h

            if img_ratio > avail_ratio:
                fit_w = avail_w
                fit_h = int(avail_w / img_ratio)
            else:
                fit_h = avail_h
                fit_w = int(avail_h * img_ratio)

            # Horizontal: always center within slot
            fit_x = slot.x + (slot.width - fit_w) // 2

            # Vertical: align by anchor
            if anchor == "top":
                fit_y = slot.y + (slot.height - avail_h) // 4  # slight top margin
            elif anchor == "bottom":
                fit_y = slot.y + slot.height - fit_h - (slot.height - avail_h) // 4
            else:  # center
                fit_y = slot.y + (slot.height - fit_h) // 2

            slide.shapes.add_picture(
                str(img_path),
                Emu(fit_x), Emu(fit_y), Emu(fit_w), Emu(fit_h),
            )
        except Exception as e:
            logger.warning("Failed to add illustration {}: {}", img_path, e)

    # ── XML Patches ──────────────────────────────────────────────

    def _patch_card_shadow(self, shape):
        """Add an outer shadow to a card shape via XML."""
        ns_a = _NSMAP["a"]
        sp_pr = shape._element.find(f".//{{{ns_a}}}spPr")
        if sp_pr is None:
            return

        shadow_color = self.design.shadow_color.lstrip("#")

        effect_lst = etree.SubElement(sp_pr, f"{{{ns_a}}}effectLst")
        outer_shdw = etree.SubElement(effect_lst, f"{{{ns_a}}}outerShdw")
        outer_shdw.set("blurRad", "381000")   # ~30pt blur — larger, softer spread
        outer_shdw.set("dist", "101600")       # ~8pt distance — more depth
        outer_shdw.set("dir", "5400000")       # Straight down
        outer_shdw.set("algn", "t")
        outer_shdw.set("rotWithShape", "0")

        srgb_clr = etree.SubElement(outer_shdw, f"{{{ns_a}}}srgbClr")
        srgb_clr.set("val", shadow_color)
        alpha = etree.SubElement(srgb_clr, f"{{{ns_a}}}alpha")
        alpha.set("val", "14000")  # 14% opacity — subtle, natural depth

    def _patch_corner_radius(self, shape, adj_val: int = 5000):
        """Set the corner radius of a rounded rectangle."""
        ns_a = _NSMAP["a"]
        prst_geom = shape._element.find(f".//{{{ns_a}}}prstGeom")
        if prst_geom is not None:
            av_lst = prst_geom.find(f"{{{ns_a}}}avLst")
            if av_lst is None:
                av_lst = etree.SubElement(prst_geom, f"{{{ns_a}}}avLst")
            # Clear existing
            for child in list(av_lst):
                av_lst.remove(child)
            gd = etree.SubElement(av_lst, f"{{{ns_a}}}gd")
            gd.set("name", "adj")
            gd.set("fmla", f"val {adj_val}")

    def _patch_svg_blip(self, pic_shape, icon_name: str):
        """Add SVG extension to a picture's blip (for modern PowerPoint)."""
        svg_str = get_icon_svg(icon_name, self.design.icon_color)
        svg_bytes = svg_str.encode("utf-8")

        # Save SVG to temp file and add as relationship
        svg_path = os.path.join(self._temp_dir, f"{icon_name}.svg")
        with open(svg_path, "wb") as f:
            f.write(svg_bytes)

        # Add SVG as a relationship on the slide part
        slide_part = pic_shape.part
        svg_rel = slide_part.relate_to_file(
            svg_path,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
        )

        # Patch the blip XML to include SVG extension
        ns_a = _NSMAP["a"]
        blip = pic_shape._element.find(f".//{{{ns_a}}}blip")
        if blip is None:
            return

        ext_lst = blip.find(f"{{{ns_a}}}extLst")
        if ext_lst is None:
            ext_lst = etree.SubElement(blip, f"{{{ns_a}}}extLst")

        # Add SVG extension
        ext = etree.SubElement(ext_lst, f"{{{ns_a}}}ext")
        ext.set("uri", "{96DAC541-7B7A-43D3-8B79-37D633B846F1}")

        ns_asvg = _NSMAP["asvg"]
        svg_blip = etree.SubElement(ext, f"{{{ns_asvg}}}svgBlip")
        svg_blip.set(f"{{{_NSMAP['r']}}}embed", svg_rel)
